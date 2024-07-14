import os
import openai
from pymongo import MongoClient
from pptx import Presentation
from docx import Document
from config import Config
from datetime import datetime

# MongoDB client and collections
client = MongoClient(Config.MONGODB_URI)
db = client[Config.MONGODB_DB]
documents_collection = db['documents']

# Initialize OpenAI API
openai.api_key = Config.OPENAI_API_KEY

def debug_log(message):
    print(f"[DEBUG] {message}")

def extract_text_from_pptx(file_path):
    debug_log(f"Extracting text from PPTX: {file_path}")
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_docx(file_path):
    debug_log(f"Extracting text from DOCX: {file_path}")
    doc = Document(file_path)
    text_runs = []
    for paragraph in doc.paragraphs:
        text_runs.append(paragraph.text)
    return "\n".join(text_runs)

def extract_text_from_txt(file_path):
    debug_log(f"Extracting text from TXT: {file_path}")
    with open(file_path, 'r') as file:
        return file.read()

def generate_question_answer(content):
    context = "Context: MongoDB Developer Days, MongoDB Atlas, MongoDB Aggregation Pipelines, and MongoDB Atlas Search"
    prompt = f"{context}\n\nBased on the following content, generate a series of questions and answers that workshop attendees may encounter:\n\n{content}"
    
    debug_log("Generating questions and answers...")
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are an assistant that generates questions and answers based on provided content."},
            {"role": "user", "content": prompt}
        ]
    )
    return response['choices'][0]['message']['content'].strip()

def save_to_mongodb(question, answer, title, summary, references):
    document = {
        'question': question,
        'title': title,
        'summary': summary,
        'answer': answer,
        'main_answer': answer,
        'references': references,
        'question_embedding': generate_embedding(question),  # Implement generate_embedding function
        'answer_embedding': generate_embedding(answer),  # Implement generate_embedding function
        'created_at': datetime.now(),
        'updated_at': datetime.now(),
        'schema_version': 2,
        'created_by': 'ai'
    }
    debug_log(f"Preparing to save document to MongoDB in database '{Config.MONGODB_DB}', collection 'documents'")
    debug_log(f"Document: {document}")
    result = documents_collection.insert_one(document)
    debug_log(f"Document inserted with ID: {result.inserted_id}")

def generate_title_and_summary(answer):
    prompt = f"Generate a concise title and summary for the following answer:\n\n{answer}\n\nProvide the title and summary in the following format:\nTitle: [Title]\nSummary: [Summary]"
    
    debug_log("Generating title and summary...")
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are an assistant that generates titles and summaries."},
            {"role": "user", "content": prompt}
        ]
    )
    
    content = response['choices'][0]['message']['content'].strip()
    title = content.split("Title: ")[1].split("Summary: ")[0].strip()
    summary = content.split("Summary: ")[1].strip()
    
    return title, summary

def generate_references(answer):
    prompt = f"Generate relevant references for the following answer, including links to the official MongoDB documentation if applicable:\n\n{answer}\n\nProvide the references in the following format:\nReferences: [References]"
    
    debug_log("Generating references...")
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are an assistant that generates references."},
            {"role": "user", "content": prompt}
        ]
    )
    
    references = response['choices'][0]['message']['content'].strip().replace("References: ", "")
    
    return references

def generate_embedding(text):
    debug_info = {}
    try:
        logger.debug(f"Generating embedding for text: {text[:50]}...")
        debug_info['openai_request'] = {
            'model': "text-embedding-ada-002",
            'input': text
        }
        response = openai.Embedding.create(model="text-embedding-ada-002", input=text)
        debug_info['openai_response'] = json.loads(json.dumps(response, default=str))
        embedding = response['data'][0]['embedding']
        debug_info['embedding_length'] = len(embedding)
        logger.debug(f"Embedding generated successfully. Length: {len(embedding)}")
        return embedding, debug_info
    except Exception as e:
        logger.error(f"Error generating embedding: {str(e)}")
        debug_info['error'] = str(e)
        debug_info['traceback'] = traceback.format_exc()
        raise

def process_files_in_directory(directory):
    debug_log(f"Processing files in directory: {directory}")
    for root, dirs, files in os.walk(directory):
        for file in files:
            debug_log(f"Processing file: {file}")
            if file.endswith('.pptx'):
                content = extract_text_from_pptx(os.path.join(root, file))
            elif file.endswith('.docx'):
                content = extract_text_from_docx(os.path.join(root, file))
            elif file.endswith('.txt'):
                content = extract_text_from_txt(os.path.join(root, file))
            else:
                debug_log(f"Skipping unsupported file: {file}")
                continue

            qa_pairs = generate_question_answer(content).split("\n\n")  # Assuming QA pairs are separated by double newlines
            for qa_pair in qa_pairs:
                if "Q:" in qa_pair and "A:" in qa_pair:
                    question = qa_pair.split("Q:")[1].split("A:")[0].strip()
                    answer = qa_pair.split("A:")[1].strip()
                    title, summary = generate_title_and_summary(answer)
                    references = generate_references(answer)
                    save_to_mongodb(question, answer, title, summary, references)

if __name__ == "__main__":
    directory = "documents/"
    process_files_in_directory(directory)
